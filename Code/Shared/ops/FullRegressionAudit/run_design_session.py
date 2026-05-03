"""GPT Iterative Design Session — UML object model, sequence diagrams, backlog."""
import json
import os
import subprocess
import time
from openai import OpenAI
from pptx import Presentation
from dotenv import load_dotenv

REPO_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", ".."))
load_dotenv(os.path.join(REPO_ROOT, "Code", ".env"))

client = OpenAI(api_key=os.environ["OPENAI_API_KEY"])


def read_file(path):
    full = os.path.join(REPO_ROOT, path)
    if not os.path.exists(full):
        return f"FILE NOT FOUND: {path}"
    with open(full, encoding="utf-8", errors="replace") as f:
        content = f.read()
    if len(content) > 30000:
        return content[:30000] + "\n...[TRUNCATED]..."
    return content


def run_git(cmd):
    try:
        result = subprocess.run(
            cmd.split(), capture_output=True, text=True,
            cwd=REPO_ROOT, encoding="utf-8", errors="replace", timeout=30,
        )
        return (result.stdout or "") + (result.stderr or "")
    except Exception as e:
        return f"ERROR: {e}"


def main():
    # 1. PPTX diagram
    prs = Presentation(os.path.join(REPO_ROOT, "Analysis/PromptngArchittecture Find CAe.pptx"))
    pptx_text = ""
    for i, slide in enumerate(prs.slides):
        pptx_text += f"\n=== Slide {i+1} ===\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        pptx_text += f"  {text}\n"

    # 2. Audit + architecture recommendation
    with open(os.path.join(REPO_ROOT, "brain/BusinessArtifacts/Audits/04_final_report.json"), encoding="utf-8") as f:
        audit = json.load(f)

    # 3. GPT final architecture
    with open(os.path.join(REPO_ROOT, "brain/BusinessArtifacts/Audits/09_gpt_final_architecture.json"), encoding="utf-8") as f:
        gpt_arch = json.load(f)

    # 4. Current code
    code_files = [
        "Code/ConversationalUX/FindCareChat/backend/main.py",
        "Code/ConversationalUX/FindCareChat/backend/application/tool_router.py",
        "Code/ConversationalUX/FindCareChat/backend/application/facades/find_care_facade.py",
        "Code/ConversationalUX/FindCareChat/backend/application/facades/evaluate_care_facade.py",
        "Code/ConversationalUX/FindCareChat/backend/domain/find_care/specialty_service.py",
        "Code/ConversationalUX/FindCareChat/backend/domain/find_care/provider_search_service.py",
        "Code/ConversationalUX/FindCareChat/backend/domain/shared/safety/safety_service.py",
        "Code/ConversationalUX/FindCareChat/backend/domain/shared/consent/consent_service.py",
        "Code/ConversationalUX/FindCareChat/backend/domain/shared/lead_capture/lead_service.py",
        "Code/ConversationalUX/FindCareChat/backend/url_guardian.py",
        "Code/ConversationalUX/FindCareChat/frontend/src/components/FindCareApp.tsx",
        "Code/ConversationalUX/FindCareChat/frontend/src/components/ChatWindow.tsx",
        "Code/Shared/prompt_system_maker.py",
    ]
    current_code = ""
    for cf in code_files:
        content = read_file(cf)
        current_code += f"\n--- {cf} ---\n{content[:15000]}\n"

    # 5. Epic structure
    with open(os.path.join(REPO_ROOT, "brain/machine_artifacts/content/agile_backlog.json"), encoding="utf-8") as f:
        backlog = json.load(f)
    epics = []
    for epic in backlog.get("epics", []):
        e = {"epic_id": epic.get("epic_id"), "name": epic.get("name"), "features": []}
        for feat in epic.get("features", []):
            e["features"].append({"feature_id": feat.get("feature_id"), "name": feat.get("name")})
        epics.append(e)
    for key, val in backlog.items():
        if isinstance(val, dict) and "epic_id" in val:
            e = {"epic_id": val.get("epic_id"), "name": val.get("name"), "features": []}
            for feat in val.get("features", []):
                e["features"].append({"feature_id": feat.get("feature_id"), "name": feat.get("name")})
            epics.append(e)

    # 6. Design.json
    design = read_file("brain/machine_artifacts/content/design.json")

    system_prompt = (
        "You are a senior software architect designing a prompt factory architecture for ChatHealthy.ai.\n\n"
        "You are in an ITERATIVE DESIGN SESSION. Multiple rounds of refinement.\n\n"
        "When you need files: {\"action\": \"request_files\", \"files\": [\"path1\"]}\n"
        "When you need git: {\"action\": \"git_command\", \"command\": \"git log ...\"}\n"
        "When producing design: {\"action\": \"design_output\", \"iteration\": N, \"data\": {...}}\n\n"
        "DELIVERABLES (build incrementally):\n\n"
        "1. UML CLASS DIAGRAM (text):\n"
        "   - Every class with attributes and methods (full signatures with params + return types)\n"
        "   - Inheritance, interfaces, composition\n"
        "   - Middleware chain\n\n"
        "2. UML SEQUENCE DIAGRAMS (text) for 7 flows:\n"
        "   a. Provider search (happy path)\n"
        "   b. Conversational question (narrative response)\n"
        "   c. Clinical trials search\n"
        "   d. Emergency/safety detection\n"
        "   e. Context carry-forward (Virginia then Richmond)\n"
        "   f. Consent trigger (5th message)\n"
        "   g. URL validation on tool response\n\n"
        "3. BACKLOG: features, stories (with size), boolean-testable requirements, test children\n\n"
        "Be rigorous. Use real class names. Show method signatures. Reference actual codebase files.\n"
        "This is a design document developers will implement from.\n\n"
        "Pattern: Interface-first. Thin superclass (just routing/dispatch). Middleware for cross-cutting.\n"
        "Subclass per epic. Lift shared logic only when proven by experience.\n\n"
        "Output JSON with action field on every response."
    )

    initial = (
        f"=== ARCHITECTURE DIAGRAM ===\n{pptx_text}\n\n"
        f"=== AGREED PATTERN ===\n{json.dumps(gpt_arch['final_architecture_recommendation'], indent=2)}\n\n"
        f"=== FINDING MAPPING ===\n{json.dumps(audit.get('architecture_recommendation', {}).get('finding_mapping', []), indent=2)}\n\n"
        f"=== EPIC STRUCTURE ===\n{json.dumps(epics, indent=2)}\n\n"
        f"=== CURRENT CODE ===\n{current_code}\n\n"
        f"=== DESIGN.JSON ===\n{design[:20000]}\n\n"
        "Begin your design. Start with the UML class diagram. Request any files you need.\n"
        "Remember: Interface-first. Thin superclass. Middleware for cross-cutting. Subclass per epic."
    )

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": initial},
    ]

    log_path = os.path.join(REPO_ROOT, "brain", "BusinessArtifacts", "Audits", "10_design_iteration_log.txt")
    log = open(log_path, "w", encoding="utf-8")
    log.write(f"Prompt Factory Design Session -- {time.strftime('%Y-%m-%d %H:%M:%S')}\n")
    log.write(f"{'='*80}\n\n")

    total_in = 0
    total_out = 0
    designs = []

    for iteration in range(150):
        print(f"\n--- Iteration {iteration + 1} ---")
        log.write(f"\n--- Iteration {iteration + 1} ---\n")

        response = client.chat.completions.create(
            model="gpt-4.1",
            temperature=0,
            messages=messages,
            response_format={"type": "json_object"},
            max_tokens=16000,
        )

        total_in += response.usage.prompt_tokens
        total_out += response.usage.completion_tokens

        gpt_msg = response.choices[0].message.content
        messages.append({"role": "assistant", "content": gpt_msg})

        try:
            action = json.loads(gpt_msg)
        except json.JSONDecodeError:
            log.write(f"  Non-JSON: {gpt_msg[:500]}\n")
            break

        action_type = action.get("action", "")
        log.write(f"  Action: {action_type}\n")

        if action_type == "request_files":
            files = action.get("files", [])
            print(f"  Requesting {len(files)} files: {[os.path.basename(f) for f in files]}")
            log.write(f"  Files: {files}\n")
            parts = []
            for fpath in files:
                content = read_file(fpath)
                parts.append(f"=== {fpath} ===\n{content}")
                log.write(f"    Served: {fpath} ({len(content)} chars)\n")
            messages.append({"role": "user", "content": "\n\n".join(parts)})

        elif action_type == "git_command":
            cmd = action.get("command", "")
            print(f"  Git: {cmd}")
            log.write(f"  Git: {cmd}\n")
            output = run_git(cmd)
            log.write(f"    Output: {len(output)} chars\n")
            messages.append({"role": "user", "content": f"Git output:\n{output}"})

        elif action_type == "design_output":
            it = action.get("iteration", iteration + 1)
            data = action.get("data", {})
            designs.append(data)
            print(f"  Design output iteration {it}")

            has_classes = bool(data.get("uml_class_diagram"))
            has_sequences = bool(data.get("uml_sequence_diagrams"))
            has_backlog = bool(data.get("backlog"))

            log.write(f"  Design v{it}: classes={has_classes} sequences={has_sequences} backlog={has_backlog}\n")

            if has_classes and has_sequences and has_backlog:
                complete = [d for d in designs if d.get("uml_class_diagram") and d.get("uml_sequence_diagrams") and d.get("backlog")]
                if len(complete) >= 2:
                    print(f"\n{'='*60}")
                    print(f"Design complete after {iteration + 1} iterations, {len(complete)} complete versions")
                    break
                messages.append({"role": "user", "content": (
                    "Good draft. Now critically review your own design:\n"
                    "1. Are all 7 sequence diagrams complete with method calls, params, return types?\n"
                    "2. Does every class have FULL method signatures (params + return types)?\n"
                    "3. Are requirements boolean-testable?\n"
                    "4. Did you miss edge cases? Error handling? Timeout flows?\n"
                    "5. Does middleware handle UMLS compliance filtering?\n"
                    "6. How does the UI know whether to render narrative vs structured results?\n"
                    "7. Where does prompt history get persisted?\n"
                    "8. How are tools registered and dispatched per epic?\n"
                    "Produce a REFINED version with action: design_output."
                )})
            else:
                missing = []
                if not has_classes: missing.append("UML class diagram")
                if not has_sequences: missing.append("sequence diagrams")
                if not has_backlog: missing.append("backlog")
                messages.append({"role": "user", "content": f"Missing: {', '.join(missing)}. Continue building."})

        else:
            print(f"  Unknown: {action_type}")
            messages.append({"role": "user", "content": "Continue with the design. Request files or produce output."})

        # Trim
        total_len = sum(len(m["content"]) for m in messages)
        if total_len > 900000:
            print(f"  Trimming ({total_len:,} chars)")
            messages = messages[:2] + messages[-30:]

    # Save
    if designs:
        final = designs[-1]
        out_path = os.path.join(REPO_ROOT, "brain", "BusinessArtifacts", "Audits", "10_prompt_factory_design.json")
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(final, f, indent=2, ensure_ascii=False)
        print(f"\nSaved: {out_path}")

    print(f"Tokens: in={total_in:,} out={total_out:,}")
    print(f"Iterations: {len(designs)} design outputs")
    log.write(f"\nTotal tokens: in={total_in:,} out={total_out:,}\n")
    log.write(f"Design iterations: {len(designs)}\n")
    log.close()


if __name__ == "__main__":
    main()
