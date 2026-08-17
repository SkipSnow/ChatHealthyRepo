"""Build the ClaudeCodeUtterances classification one-pager (PowerPoint).

Lives beside the classifier that produces tone / utterance_class values.
Reads a cached Mongo aggregate JSON and logo from `_oneshots/`, writes the
slide there as well.
"""
from __future__ import annotations

import json
from datetime import date
from pathlib import Path

from chathealthy_lib.exceptions import ChatHealthyException
from chathealthy_lib.logging_service import ChatHealthyLoggingService

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

_log = ChatHealthyLoggingService()

_HERE = Path(__file__).resolve().parent
ROOT = _HERE.parents[1]
ONESHOTS = ROOT / "_oneshots"
DATA = json.loads((ONESHOTS / "utterance_ops_report_data.json").read_text(encoding="utf-8"))

logo_tif = ONESHOTS / "logo.tif"
logo_png = ONESHOTS / "logo_for_pptx.png"
Image.open(logo_tif).save(logo_png, format="PNG")

SLIDE_W = Inches(8)
SLIDE_H = Inches(11.5)

NAVY = RGBColor(0x0F, 0x2A, 0x44)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x64, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BODY_BG = RGBColor(240, 248, 252)
LIGHT_GREY = RGBColor(236, 238, 241)
FOOTER_BG = LIGHT_GREY
HEADER_ROW_BG = LIGHT_GREY
FRAME_BORDER = RGBColor(0x6E, 0x7F, 0x8F)
LINE = RGBColor(0xB8, 0xC2, 0xCC)
ROW_BAND = RGBColor(0xF4, 0xF8, 0xFB)
FONT = "Calibri"

total = int(DATA["total"])


def _usable(rows: list[dict]) -> list[dict]:
    return [r for r in rows if str(r.get("value")) != "(missing)"]


tones = _usable(sorted(DATA["tones"], key=lambda r: (-int(r["n"]), str(r["value"]))))
classes = _usable(sorted(DATA["classes"], key=lambda r: (-int(r["n"]), str(r["value"]))))

HUMAN_ROLES = {"user"}
BOT_ROLES = {"assistant", "agent", "service"}
human_total = sum(
    int(r["n"])
    for r in DATA.get("roles", [])
    if str(r.get("_id") or "").lower() in HUMAN_ROLES
)
bot_total = sum(
    int(r["n"])
    for r in DATA.get("roles", [])
    if str(r.get("_id") or "").lower() in BOT_ROLES
)
ACTOR_DENOM = {"Human actor": human_total, "Code Bot": bot_total}


def display_name(value: str) -> str:
    raw = str(value)
    special = {
        "report": "Status query & result",
        "matter_of_fact": "Matter of fact",
        "source_material": "Source material",
        "unclassified": "Unclassified",
    }
    if raw in special:
        return special[raw]
    return raw.replace("_", " ").capitalize()


def professional_samples(items: list[dict], limit: int = 6) -> list[dict]:
    picked, seen = [], set()
    for item in items:
        content = (item.get("content") or "").strip()
        if not content or content in seen:
            continue
        if not content.strip("/?.,!;:-_|\\") or len(content) < 2:
            continue
        seen.add(content)
        picked.append(item)
        if len(picked) >= limit:
            break
    return picked


samples_tone = professional_samples(DATA["samples_tone_unclassified"])
samples_class = professional_samples(DATA["samples_class_unclassified"])

# Four panels with clear body-blue gutters
GUTTER_X = Inches(0.45)
MARGIN_X = Inches(0.18)
COL_W = Inches((8.0 - 2 * 0.18 - 0.45) / 2)
LEFT_COL_X = MARGIN_X
RIGHT_COL_X = MARGIN_X + COL_W + GUTTER_X

FOOTER_H = Inches(0.30)
FOOTER_TOP = SLIDE_H - FOOTER_H
GUTTER_Y = Inches(0.30)
STATS_TOP = Inches(0.68)
SAMPLES_TOP = Inches(8.90)
STATS_H = SAMPLES_TOP - STATS_TOP - GUTTER_Y
FRAME_PAD = Inches(0.10)

F_TITLE = 15
F_SUBTITLE = 9.5
F_SECTION = 12
F_NOTE = 6.5
F_HDR = 9
F_SUMMARY = 8
F_ACTOR = 7
F_SAMPLE_TITLE = 10
F_SAMPLE = 7.5
F_FOOTER = 8

# Exact line geometry (textbox rows — no PPT table height surprises)
H_LINE_MAIN = Inches(0.140)
H_LINE_SUB = Inches(0.118)
H_ITEM_PAD_BOTTOM = Inches(0.036)
H_ITEM = H_LINE_MAIN + 2 * H_LINE_SUB + H_ITEM_PAD_BOTTOM
GAP_BETWEEN_ITEMS = Inches(0.115)  # clearly larger than within-group (0)
H_HDR_BAR = Inches(0.38)

W0, W1, W2 = 0.56, 0.24, 0.20


def pct(n: int, denom: int | None = None) -> str:
    base = total if denom is None else denom
    if base <= 0:
        return "0.0%"
    return f"{(100.0 * n / base):.1f}%"


def set_run(run, text, size_pt, bold=False, italic=False, color=INK, font_name=FONT):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


def _clear_effects(shape) -> None:
    """Prevent theme/export drop-shadows from cutting through text."""
    spPr = shape._element.spPr
    for child in list(spPr):
        tag = child.tag
        if tag.endswith("}effectLst") or tag.endswith("}scene3d") or tag.endswith("}sp3d"):
            spPr.remove(child)
    etree.SubElement(spPr, qn("a:effectLst"))


def add_rect(slide, left, top, width, height, fill: RGBColor, line: RGBColor | None = None, line_pt: float = 1.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_pt)
    _clear_effects(shape)
    return shape


def add_cell_text(slide, left, top, width, height, text, size_pt, *, bold=False, color=INK, pad_pt=2):
    from pptx.enum.text import MSO_AUTO_SIZE

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(pad_pt)
    p.space_after = Pt(0)
    set_run(p.add_run(), text, size_pt, bold=bold, color=color)
    return box


def add_header_bar(slide, left, top, width):
    add_rect(slide, left, top, width, H_HDR_BAR, HEADER_ROW_BG, line=None)
    xs = [left, left + width * W0, left + width * W0 + width * W1]
    ws = [width * W0, width * W1, width * W2]
    text_top = top + Inches(0.08)
    text_h = H_HDR_BAR - Inches(0.12)
    for x, w, label in zip(xs, ws, ("Classification", "Count", "%")):
        add_cell_text(
            slide,
            x + Inches(0.04),
            text_top,
            w - Inches(0.04),
            text_h,
            label,
            F_HDR,
            bold=True,
            color=NAVY,
            pad_pt=0,
        )
    return top + H_HDR_BAR


def add_item_block(slide, left, top, width, row: dict, band: RGBColor):
    """Main + 2 actor lines, tightly stacked; returns y after the block."""
    name = display_name(row["value"])
    n = int(row["n"])
    human = int(row.get("human", 0))
    bot = int(row.get("bot", 0))
    actors = sorted(
        (("Code Bot", bot), ("Human actor", human)),
        key=lambda pair: -pair[1],
    )

    # Single fill for the whole group — no internal borders that clip text
    add_rect(slide, left, top, width, H_ITEM, band, line=None)

    xs = [left, left + width * W0, left + width * W0 + width * W1]
    ws = [width * W0, width * W1, width * W2]

    def write_line(y, height, c0, c1, c2, *, bold, color, size):
        add_cell_text(slide, xs[0] + Inches(0.04), y, ws[0] - Inches(0.04), height, c0, size, bold=bold, color=color, pad_pt=3)
        add_cell_text(slide, xs[1] + Inches(0.02), y, ws[1] - Inches(0.02), height, c1, size, bold=bold, color=color, pad_pt=3)
        add_cell_text(slide, xs[2] + Inches(0.02), y, ws[2] - Inches(0.02), height, c2, size, bold=bold, color=color, pad_pt=3)

    y = top
    write_line(y, H_LINE_MAIN, name, f"{n:,}", pct(n), bold=True, color=INK, size=F_SUMMARY)
    y += H_LINE_MAIN
    for actor_label, actor_n in actors:
        write_line(
            y,
            H_LINE_SUB,
            f"  {actor_label}",
            f"{actor_n:,}",
            pct(actor_n, ACTOR_DENOM[actor_label]),
            bold=False,
            color=MUTED,
            size=F_ACTOR,
        )
        y += H_LINE_SUB
    return top + H_ITEM


def write_stat_column(slide, left, title: str, rows: list[dict]):
    add_rect(slide, left, STATS_TOP, COL_W, STATS_H, WHITE, line=FRAME_BORDER, line_pt=1.25)

    inner_left = left + FRAME_PAD
    inner_w = COL_W - 2 * FRAME_PAD
    y = STATS_TOP + Inches(0.06)

    title_box = slide.shapes.add_textbox(inner_left, y, inner_w, Inches(0.22))
    p = title_box.text_frame.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    set_run(p.add_run(), title, F_SECTION, bold=True, color=NAVY)
    y += Inches(0.22)

    note_box = slide.shapes.add_textbox(inner_left, y, inner_w, Inches(0.14))
    p = note_box.text_frame.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    set_run(p.add_run(), "Actor lines: % of that actor’s utterances", F_NOTE, italic=True, color=MUTED)
    y += Inches(0.18)

    y = add_header_bar(slide, inner_left, y, inner_w)
    y += Inches(0.10)  # clear separation so header never overlaps first item

    needed = len(rows) * float(H_ITEM) + max(0, len(rows) - 1) * float(GAP_BETWEEN_ITEMS)
    avail = float(STATS_TOP + STATS_H - Inches(0.06)) - float(y)
    if needed > avail:
        raise ChatHealthyException(
            mode="report_column_overflow",
            component="build_utterance_operational_report",
            message=(
                f"'{title}' needs {needed/914400:.2f}\" but only "
                f"{avail/914400:.2f}\" remains ({len(rows)} groups)"
            ),
        )

    for idx, row in enumerate(rows):
        band = ROW_BAND if idx % 2 == 0 else WHITE
        y = add_item_block(slide, inner_left, y, inner_w, row, band)
        if idx < len(rows) - 1:
            y += GAP_BETWEEN_ITEMS  # white frame shows through = bigger gap between items


def write_sample_column(slide, left, title: str, items: list[dict]):
    panel_h = FOOTER_TOP - SAMPLES_TOP - Inches(0.04)
    add_rect(slide, left, SAMPLES_TOP, COL_W, panel_h, WHITE, line=FRAME_BORDER, line_pt=1.25)

    box = slide.shapes.add_textbox(
        left + FRAME_PAD,
        SAMPLES_TOP + Inches(0.10),
        COL_W - 2 * FRAME_PAD,
        panel_h - Inches(0.14),
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after = Pt(6)
    set_run(p.add_run(), title, F_SAMPLE_TITLE, bold=True, color=NAVY)
    for item in items:
        actor = item.get("actor") or "—"
        content = (item.get("content") or "").replace("\n", " ").strip()
        if len(content) > 78:
            content = content[:75] + "…"
        para = tf.add_paragraph()
        para.space_before = Pt(3)
        para.space_after = Pt(0)
        set_run(para.add_run(), f"{actor}:  “{content}”", F_SAMPLE, color=MUTED)


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BODY_BG)
add_rect(slide, Inches(0), FOOTER_TOP, SLIDE_W, FOOTER_H, FOOTER_BG)

logo_h = Inches(0.34)
slide.shapes.add_picture(str(logo_png), Inches(0.18), Inches(0.12), height=logo_h)

title_box = slide.shapes.add_textbox(Inches(2.15), Inches(0.08), Inches(5.65), Inches(0.50))
tf = title_box.text_frame
tf.word_wrap = True
p0 = tf.paragraphs[0]
p0.space_before = Pt(0)
p0.space_after = Pt(0)
set_run(p0.add_run(), "Conversation Classification Summary", F_TITLE, bold=True, color=NAVY)
p1 = tf.add_paragraph()
p1.space_before = Pt(2)
p1.space_after = Pt(0)
set_run(
    p1.add_run(),
    f"ChatHealthy.ai engineering dialogue archive  ·  {total:,} utterances  ·  Tone × Class",
    F_SUBTITLE,
    color=MUTED,
)

write_stat_column(slide, LEFT_COL_X, "Semantic classification", classes)
write_stat_column(slide, RIGHT_COL_X, "Tone classification", tones)
write_sample_column(slide, LEFT_COL_X, "Unclassified examples — Class", samples_class)
write_sample_column(slide, RIGHT_COL_X, "Unclassified examples — Tone", samples_tone)

footer_y = FOOTER_TOP + Inches(0.05)
for left, width, text, align, bold in (
    (Inches(0.18), Inches(2.2), date.today().strftime("%B %d, %Y"), PP_ALIGN.LEFT, False),
    (Inches(2.0), Inches(4.0), "ChatBot conversation summary", PP_ALIGN.CENTER, True),
    (Inches(5.55), Inches(2.25), "© ChatHealthy.ai", PP_ALIGN.RIGHT, False),
):
    box = slide.shapes.add_textbox(left, footer_y, width, Inches(0.20))
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text, F_FOOTER, bold=bold, color=MUTED if not bold else NAVY)

out_pptx = ONESHOTS / "ClaudeCodeUtterancesClassificationOpsReport.pptx"
prs.save(str(out_pptx))
_log.info("Wrote %s", out_pptx)
_log.info("groups class=%d tone=%d", len(classes), len(tones))
